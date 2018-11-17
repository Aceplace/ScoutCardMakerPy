import os

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

from misc.scoutcardmakerexceptions import ScoutCardMakerException
from offensiveformation.formation import Formation

#Constants for wide view
CENTER_X_POS = Cm(13.0)
CENTER_Y_POS = Cm(11.5)
HORIZONTAL_COORDINATE_SIZE = Cm(0.23)
VERTICAL_COORDINATE_SIZE = Cm(0.55)
PLAYER_WIDTH = Cm(0.65)
PLAYER_HEIGHT = Cm(0.5)
HASH_SIZE = Cm(0.1)
TITLE_LEFT = 106531
TITLE_TOP = 132963
TITLE_WIDTH = 8904303
TITLE_HEIGHT = 369332
DEFENDER_WIDTH = Pt(28)
DEFENDER_HEIGHT = Pt(28)
#constants for tight view
TIGHT_CENTER_X_POS = Cm(13.0)
TIGHT_CENTER_Y_POS = Cm(11.0)
TIGHT_HORIZONTAL_COORDINATE_SIZE = Cm(0.46)
TIGHT_VERTICAL_COORDINATE_SIZE = Cm(1.05)
TIGHT_PLAYER_WIDTH = Cm(1.3)
TIGHT_PLAYER_HEIGHT = Cm(1.0)
TIGHT_DEFENDER_WIDTH = Pt(40)
TIGHT_DEFENDER_HEIGHT = Pt(40)
#Derived Constant Values
LEFT_SIDELINE = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 53
RIGHT_SIDELINE = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 53
LEFT_HASH = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 18
RIGHT_HASH = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 18
LEFT_TOP_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 35
LEFT_BOTTOM_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 39
RIGHT_TUP_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 35
RIGHT_BOTTOM_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 39
FIVE_YARDS = VERTICAL_COORDINATE_SIZE * 5


def player_coordinates_to_powerpoint(player_x, player_y, is_tight_view=False):
    if not is_tight_view:
        return (CENTER_X_POS + player_x * HORIZONTAL_COORDINATE_SIZE, CENTER_Y_POS + player_y * VERTICAL_COORDINATE_SIZE)
    else:
        return (TIGHT_CENTER_X_POS + player_x * TIGHT_HORIZONTAL_COORDINATE_SIZE, TIGHT_CENTER_Y_POS + player_y * TIGHT_VERTICAL_COORDINATE_SIZE)


def export_to_powerpoint(output_filename, plays, formation_library, defense_library):

    presentation = Presentation()

    #do wide versions first
    for play in plays:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]} {play["Defense"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        if play['Note']:
            text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
            text_box.text_frame.text = f'{play["Note"]}'
            text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_wide_formation_and_defense_slide(play, slide, formation_library, defense_library)

    # do tight versions after
    for play in plays:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]} {play["Defense"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        if play['Note']:
            text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
            text_box.text_frame.text = f'{play["Note"]}'
            text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_tight_formation_and_defense_slide(play, slide, formation_library, defense_library)

    presentation.save(output_filename)

    write_missing_formations_and_defenses_to_text_file(output_filename, plays, formation_library, defense_library)


def add_wide_formation_and_defense_slide(play, slide, formation_library, defense_library):
    #draw sideline
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)

    for num in range(-3, 3):
        #draw lines that go accross field at 5 yard intervals
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num)
        #draw hash marks
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        #draw ticks for the top of numbers and bottom of numbers
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)

    formation_name = play['Card Maker Formation']
    defense_name = play['Card Maker Defense']

    formation = None
    try:
        if formation_name:
            formation = formation_library.get_composite_formation(formation_name)
            for label, player in formation.players.items():
                x, y = player_coordinates_to_powerpoint(player.x, player.y)
                shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - PLAYER_WIDTH / 2, y - PLAYER_HEIGHT / 2, PLAYER_WIDTH, PLAYER_HEIGHT)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(0, 0, 0)
                shape.line.width = Pt(1.0)
                shape.text_frame.text = player.label if len(player.label) != 2 else player.label[1]
                shape.text_frame.paragraphs[0].font.size = Pt(12)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        if formation and defense_name:
            composite_defense = defense_library.get_composite_defense(defense_name)
            placed_defense = composite_defense.get_placed_defenders(formation)
            for (label, defender_x, defender_y) in placed_defense:
                x, y = player_coordinates_to_powerpoint(defender_x, defender_y * -1)
                text_box = slide.shapes.add_textbox(x - DEFENDER_WIDTH / 2, y - DEFENDER_HEIGHT / 2, DEFENDER_WIDTH, DEFENDER_HEIGHT)
                text_box.text_frame.text = label
                text_box.text_frame.paragraphs[0].font.size = Pt(24)
    except ScoutCardMakerException:
        pass #the only error that can a formation or defendse  doesn't exist
        #if so, we ignore it and go to the next card

def add_tight_formation_and_defense_slide(play, slide, formation_library, defense_library):
    formation_name = play['Card Maker Formation']
    defense_name = play['Card Maker Defense']

    formation = None
    try:
        if formation_name:
            formation = formation_library.get_composite_formation(formation_name)
            for label, player in formation.players.items():
                if player.x < -20 or player.x > 20:
                    continue
                x, y = player_coordinates_to_powerpoint(player.x, player.y, is_tight_view=True)
                shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - TIGHT_PLAYER_WIDTH / 2, y - TIGHT_PLAYER_HEIGHT / 2,
                                               TIGHT_PLAYER_WIDTH, TIGHT_PLAYER_HEIGHT)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(0, 0, 0)
                shape.line.width = Pt(1.0)
                shape.text_frame.text = player.label if len(player.label) != 2 else player.label[1]
                shape.text_frame.paragraphs[0].font.size = Pt(24)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        if formation and defense_name:
            composite_defense = defense_library.get_composite_defense(defense_name)
            placed_defense = composite_defense.get_placed_defenders(formation)
            for (label, defender_x, defender_y) in placed_defense:
                if defender_x < -20 or defender_x > 20 or defender_y > 8:
                    continue
                x, y = player_coordinates_to_powerpoint(defender_x, defender_y * -1, is_tight_view=True)
                text_box = slide.shapes.add_textbox(x - TIGHT_DEFENDER_WIDTH / 2, y - TIGHT_DEFENDER_HEIGHT / 2, TIGHT_DEFENDER_WIDTH,
                                                    TIGHT_DEFENDER_HEIGHT)
                text_box.text_frame.text = label
                text_box.text_frame.paragraphs[0].font.size = Pt(36)
    except ScoutCardMakerException:
        pass #the only error that can a formation or defendse  doesn't exist
        #if so, we ignore it and go to the next card

def write_missing_formations_and_defenses_to_text_file(presentation_filename, plays, formation_library, defense_library):

    formations_listed_in_script = [play['Card Maker Formation'] for play in plays]
    parsed_formations = []
    for composite_formation in formations_listed_in_script:
        split_formations = composite_formation.split()
        for formation in split_formations:
            if formation != 'LT' and formation != 'RT' and formation not in parsed_formations:
                parsed_formations.append(formation)

    missing_formations = []
    for parsed_formation in parsed_formations:
        if not formation_library.does_formation_exist(parsed_formation + ' RT'):
            missing_formations.append(parsed_formation)


    defenses_listed_in_script = [play['Card Maker Defense'] for play in plays]
    parsed_defenses = []
    for composite_defense in defenses_listed_in_script:
        split_defenses = composite_defense.split()
        for defense in split_defenses:
            if defense not in parsed_formations:
                parsed_defenses.append(defense)

    missing_defenses = []
    for parsed_defense in parsed_defenses:
        if not defense_library.does_defense_exist(parsed_defense):
            missing_defenses.append(parsed_defense)

    if missing_formations or missing_defenses:
        try:
            presentation_directory = os.path.dirname(presentation_filename)
            missing_stuff_filename = os.path.join(presentation_directory, 'missingFormationsandDefense.txt')
            with open(missing_stuff_filename, 'w') as missing_stuff_file:
                if missing_formations:
                    missing_stuff_file.write('Missing Formations\n-------------------\n')
                    missing_stuff_file.write('\n'.join(missing_formations))
                    missing_stuff_file.write('\n\n')
                if missing_defenses:
                    missing_stuff_file.write('Missing Defenses\n-------------------\n')
                    missing_stuff_file.write('\n'.join(missing_defenses))
                    missing_stuff_file.write('\n\n')
        except IOError as e:
            raise ScoutCardMakerException('Couldn\'t write missing formations/defenses to text file')

















