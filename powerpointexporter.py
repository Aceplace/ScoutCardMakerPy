from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.shapes.connector import Connector
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from formation import Formation

#Constant
CENTER_X_POS = Cm(13.2)
CENTER_Y_POS = Cm(11.5)
HORIZONTAL_COORDINATE_SIZE = Cm(0.23)
VERTICAL_COORDINATE_SIZE = Cm(0.55)
PLAYER_WIDTH = Cm(0.55)
PLAYER_HEIGHT = Cm(0.45)
LABEL_FONT_SIZE = Cm(12)
HASH_SIZE = Cm(0.1)
TITLE_LEFT = 106531
TITLE_TOP = 132963
TITLE_WIDTH = 8904303
TITLE_HEIGHT = 369332
#Derived Constant Values
LEFT_SIDELINE = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 53
RIGHT_SIDELINE = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 53
LEFT_HASH = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 18
RIGHT_HASH = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 18
LEFT_TOP_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 35
LEFT_BOTTOM_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 39
RIGHT_TUP_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 35
RIGHT_BOTTOM_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 39
FIVE_YARDS = VERTICAL_COORDINATE_SIZE * 5


def player_coordinates_to_powerpoint(player):
    return (CENTER_X_POS + player.x * HORIZONTAL_COORDINATE_SIZE, CENTER_Y_POS + player.y * VERTICAL_COORDINATE_SIZE)

def export_to_powerpoint(output_filename, plays):

    presentation = Presentation()

    for play in plays:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_formation_to_slide(play["Formation"], slide)
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = play["Number"] + ' ' + play["Hash"] + ' ' + play["FormationName"]+ ' ' + play["Play"]
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
    
    presentation.save(output_filename)

def add_formation_to_slide(formation, slide):
    #draw sideline
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)

    for num in range(-3, 3):
        #draw lines that go accross field at 5 yard intervals
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num)
        #draw hash marks
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        #draw ticks for the top of numbers and bottom of numbers
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)

    #draw players
    for label, player in formation.players.items():
        x, y = player_coordinates_to_powerpoint(player)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - PLAYER_WIDTH / 2, y - PLAYER_HEIGHT / 2, PLAYER_WIDTH, PLAYER_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(1.0)
        if len(player.label) != 3 and player.label != "W":
            shape.text_frame.text = player.label
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)



formations = [Formation(), Formation()]
formations[1].flip_formation()


play1 = {
    "Number":"1",
    "Hash":"L",
    "FormationName":"Blah",
    "Play":"Kah",
    "Formation":formations[0]
}

play2 = {
    "Number":"2",
    "Hash":"C",
    "FormationName":"Blsah",
    "Play":"Ksah",
    "Formation":formations[1]
}

plays = [play1, play2]

export_to_powerpoint("output.pptx", plays)
